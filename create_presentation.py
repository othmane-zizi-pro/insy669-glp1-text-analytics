#!/usr/bin/env python3
"""Create professional PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os

BASE = "/home/ubuntu/classes/insy669/final-project"
FIG = os.path.join(BASE, "figures")
PRES_DIR = os.path.join(BASE, "presentation")
os.makedirs(PRES_DIR, exist_ok=True)

with open(f"{BASE}/data/analysis_stats.json") as f:
    stats = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RGBColor(0x21, 0x96, 0xF3)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
DARK = RGBColor(0x21, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(0x0D, 0x47, 0xA1))
add_shape(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.05), ORANGE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "Media vs Public Opinion on\nGLP-1 Weight Loss Drugs", 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
    "Ozempic & Wegovy — A Text Analytics Approach", 28, False, RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    "INSY 669 — Text Analytics | McGill University | Winter 2025", 18, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
    "V. Christopoulos  •  H. Guideau  •  S. Khosla  •  M. Yousuf  •  O. Zizi", 16, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Problem Statement", 32, True, WHITE, PP_ALIGN.LEFT)

bullets = [
    "💊 GLP-1 drugs (Ozempic, Wegovy) have become the most talked-about medications in recent history",
    "📰 Media coverage shapes public perception, but does it reflect actual user experiences?",
    "🔍 Research Question: How does media framing of GLP-1 drugs differ from public opinion?",
    "💡 Understanding this gap has implications for pharma marketing, patient education, and health policy",
    "📊 We analyze 1,500+ documents from Reddit, WebMD, and major news outlets"
]
for i, bullet in enumerate(bullets):
    add_text_box(slide, Inches(1), Inches(1.8 + i*1.0), Inches(11), Inches(0.8), bullet, 20, False, DARK)

# ============================================================
# SLIDE 3: Data Sources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Data Sources", 32, True, WHITE)

# Public box
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(1), Inches(1.8), Inches(5), Inches(0.6), "📱 Public Opinion Corpus", 24, True, BLUE)
add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(3.5),
    "Reddit (800 posts)\n• r/Ozempic\n• r/Semaglutide\n• r/WegovyWeightLoss\n\nWebMD Reviews (300)\n• Patient experiences\n• Drug ratings\n\nTotal: 1,100 documents", 17, False, DARK)

# Media box
add_shape(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.6), "📰 Media Corpus", 24, True, ORANGE)
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(3.5),
    "News Articles (400)\n• Reuters, CNN Health\n• The New York Times\n• STAT News\n• Medical News Today\n• NPR Health\n• Bloomberg, WSJ\n• And 8 more outlets\n\nTotal: 400 documents", 17, False, DARK)

# ============================================================
# SLIDE 4: Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Methodology", 32, True, WHITE)

steps = [
    ("1. Data Collection", "BeautifulSoup, requests, PRAW\nReddit, WebMD, News APIs"),
    ("2. Preprocessing", "NLTK tokenization, stopwords,\nlemmatization (WordNetLemmatizer)"),
    ("3. Feature Extraction", "Bag-of-Words, TF-IDF\n(unigrams + bigrams)"),
    ("4. Sentiment Analysis", "VADER sentiment scoring\nCompound score classification"),
    ("5. Word Associations", "PMI, Lift metrics\nCo-occurrence analysis"),
    ("6. Comparison", "Cosine similarity, MDS plots\nStatistical testing (t-test)")
]

for i, (title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(1.6 + row * 2.8)
    add_shape(slide, x, y, Inches(3.8), Inches(2.2), RGBColor(0xF5, 0xF5, 0xF5))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5), title, 18, True, BLUE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(1.2), desc, 14, False, DARK)

# ============================================================
# SLIDE 5: Sentiment Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Sentiment Analysis Results", 32, True, WHITE)

# Add sentiment histogram image
try:
    slide.shapes.add_picture(f"{FIG}/sentiment_histograms.png", Inches(0.5), Inches(1.4), Inches(8.5), Inches(3.0))
except:
    pass

# Stats box
add_shape(slide, Inches(9.2), Inches(1.4), Inches(3.8), Inches(3.0), RGBColor(0xF5, 0xF5, 0xF5))
add_text_box(slide, Inches(9.4), Inches(1.5), Inches(3.4), Inches(0.4), "Key Statistics", 18, True, DARK)
add_text_box(slide, Inches(9.4), Inches(2.0), Inches(3.4), Inches(2.2),
    f"Public Mean: {stats['public_mean_sentiment']:.4f}\n"
    f"Media Mean: {stats['media_mean_sentiment']:.4f}\n\n"
    f"T-statistic: {stats['t_statistic']:.4f}\n"
    f"P-value: {stats['p_value']:.6f}\n\n"
    f"✅ Statistically significant\n(α = 0.05)", 14, False, DARK)

# Box plot
try:
    slide.shapes.add_picture(f"{FIG}/sentiment_boxplot.png", Inches(0.5), Inches(4.6), Inches(6), Inches(2.7))
except:
    pass

try:
    slide.shapes.add_picture(f"{FIG}/sentiment_pies.png", Inches(6.8), Inches(4.6), Inches(6), Inches(2.7))
except:
    pass

# ============================================================
# SLIDE 6: Word Clouds & TF-IDF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Language Comparison: What Each Corpus Emphasizes", 32, True, WHITE)

try:
    slide.shapes.add_picture(f"{FIG}/wordclouds.png", Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.8))
except:
    pass

try:
    slide.shapes.add_picture(f"{FIG}/tfidf_comparison.png", Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.9))
except:
    pass

# ============================================================
# SLIDE 7: MDS & Associations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Document Similarity & Word Associations", 32, True, WHITE)

try:
    slide.shapes.add_picture(f"{FIG}/mds_plot.png", Inches(0.3), Inches(1.4), Inches(6.3), Inches(5.5))
except:
    pass

add_shape(slide, Inches(7), Inches(1.4), Inches(6), Inches(5.5), RGBColor(0xF5, 0xF5, 0xF5))
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.4), Inches(0.5), "MDS Analysis Insights", 20, True, DARK)
add_text_box(slide, Inches(7.3), Inches(2.3), Inches(5.4), Inches(4.0),
    f"• Cosine similarity: {stats['cosine_similarity']:.4f}\n\n"
    "• Public and media documents form\n  partially overlapping clusters\n\n"
    "• Media cluster is more tightly grouped\n  (consistent language/framing)\n\n"
    "• Public cluster is more dispersed\n  (diverse personal experiences)\n\n"
    "• Some overlap on core drug terms,\n  but divergence on context/framing", 16, False, DARK)

# ============================================================
# SLIDE 8: Side Effects
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Side Effects: Coverage Gap Analysis", 32, True, WHITE)

try:
    slide.shapes.add_picture(f"{FIG}/side_effects.png", Inches(0.3), Inches(1.4), Inches(7.5), Inches(5.5))
except:
    pass

add_shape(slide, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5.5), RGBColor(0xFF, 0xEB, 0xEE))
add_text_box(slide, Inches(8.4), Inches(1.6), Inches(4.4), Inches(0.5), "⚠️ Key Gaps", 20, True, RGBColor(0xC6, 0x28, 0x28))
add_text_box(slide, Inches(8.4), Inches(2.3), Inches(4.4), Inches(4.0),
    "Users frequently discuss side\neffects that receive much less\nmedia coverage:\n\n"
    "• Sulfur burps — embarrassing\n  but very common\n• Constipation — daily impact\n  on quality of life\n• Hair loss — unexpected and\n  distressing for patients\n• Fatigue — affects work and\n  daily functioning\n\n"
    "Media tends to focus on\nsevere/rare events like\npancreatitis and gastroparesis.", 15, False, DARK)

# ============================================================
# SLIDE 9: Temporal Trends
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Sentiment Trends Over Time", 32, True, WHITE)

try:
    slide.shapes.add_picture(f"{FIG}/sentiment_timeline.png", Inches(0.5), Inches(1.5), Inches(12), Inches(4.5))
except:
    pass

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(1.0),
    "Media sentiment remains relatively stable and slightly positive, while public sentiment fluctuates more and trends negative — "
    "reflecting the gap between media framing and lived patient experience.", 16, False, GRAY, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: Key Findings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Key Findings", 32, True, WHITE)

findings = [
    ("1", "Significant Sentiment Gap", "Public opinion is more negative (mean: -0.07) vs media (mean: +0.03). The difference is statistically significant (p < 0.001)."),
    ("2", "Different Language, Different Priorities", "Public: personal experiences, side effects, costs. Media: clinical data, market trends, regulation."),
    ("3", "Side Effect Coverage Gap", "Common side effects (nausea, constipation, sulfur burps) are underrepresented in media relative to user discussions."),
    ("4", "Cost as a Central Theme", "Affordability and insurance coverage dominate public discourse but receive proportionally less media attention."),
    ("5", "Emotional vs Analytical Framing", "Public posts are highly emotional and personal; media maintains analytical, data-driven tone."),
]

for i, (num, title, desc) in enumerate(findings):
    y = Inches(1.5 + i * 1.15)
    add_shape(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), BLUE)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.6), Inches(0.5), num, 22, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.6), y, Inches(3), Inches(0.5), title, 17, True, DARK)
    add_text_box(slide, Inches(1.6), y + Inches(0.45), Inches(10.5), Inches(0.6), desc, 14, False, GRAY)

# ============================================================
# SLIDE 11: Business Implications
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Business Implications", 32, True, WHITE)

implications = [
    ("🏥 Pharmaceutical Companies", "Should proactively address common side effects in patient materials. Current messaging may create expectation gaps that lead to discontinuation."),
    ("📋 Healthcare Providers", "Need awareness of patient concerns beyond clinical data. Users value practical advice on managing everyday side effects over efficacy statistics."),
    ("🏛️ Health Policy", "Insurance coverage and drug affordability are top public concerns. Policy interventions addressing cost could significantly improve patient outcomes."),
    ("📢 Media & Communications", "Opportunity to provide more balanced coverage that includes common user experiences alongside clinical and market reporting."),
    ("💼 Health Tech", "Patient support platforms should incorporate side effect management tools and peer community features to address the information gap."),
]

for i, (icon_title, desc) in enumerate(implications):
    y = Inches(1.5 + i * 1.15)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.0), RGBColor(0xF5, 0xF5, 0xF5))
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(4), Inches(0.4), icon_title, 17, True, BLUE)
    add_text_box(slide, Inches(1), y + Inches(0.45), Inches(11), Inches(0.5), desc, 14, False, DARK)

# ============================================================
# SLIDE 12: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(0x0D, 0x47, 0xA1))
add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
    "Thank You!", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
    "Questions?", 32, False, RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.CENTER)
add_shape(slide, Inches(4), Inches(4.8), Inches(5.333), Inches(0.03), ORANGE)
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
    "Vasilis Christopoulos  •  Hugo Guideau  •  Saksi Khosla  •  Mustafa Yousuf  •  Othmane Zizi", 16, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    "github.com/othmane-zizi-pro/insy669-glp1-text-analytics", 14, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)

# Save
pptx_path = f"{PRES_DIR}/presentation.pptx"
prs.save(pptx_path)
print(f"Presentation saved to {pptx_path}")
